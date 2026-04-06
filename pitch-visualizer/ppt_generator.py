from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import os

def create_ppt(images, captions, prompts, output_path="storyboard.pptx"):
    """
    Generate a professional PPTX storyboard from the provided images and metadata.
    """
    prs = Presentation()
    
    # Use wide aspect ratio (16:9) if possible, but default is 4:3
    # prs.slide_width = Inches(13.33)
    # prs.slide_height = Inches(7.5)

    for i in range(len(images)):
        # Select slide layout 5 (Title Only / Blank-ish)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add a title to the slide
        title_shape = slide.shapes.title
        title_shape.text = f"Scene {i+1}"

        # Ensure image path is valid
        img_path = Path(images[i])
        if not img_path.exists():
            print(f"[PPT Error] Image not found: {img_path}")
            continue

        # Add the picture (centered roughly)
        # 1 inch left, 1.2 inches top, 8 inches wide approx.
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.2), width=Inches(9))

        # Add caption text box at the bottom
        caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        tf = caption_box.text_frame
        tf.word_wrap = True
        tf.text = f"Narrative: {captions[i]}"

        # Add prompt text box (smaller font / secondary)
        if prompts and i < len(prompts):
            prompt_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.5))
            pf = prompt_box.text_frame
            pf.text = f"AI Prompt: {prompts[i][:150]}..." # Truncate for slide safety

    prs.save(output_path)
    return str(output_path)
